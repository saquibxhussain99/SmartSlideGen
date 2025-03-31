import os
import json
import uuid
import re
from PIL import Image
from dotenv import load_dotenv
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from prompt_templates import GENERATE_SLIDE_CONTENT_TEMPLATE

# === SETTINGS ===
MAX_BULLETS_PER_SLIDE = 6
MAX_IMAGE_WIDTH_INCHES = 5.5
MAX_IMAGE_HEIGHT_INCHES = 4.0
TEXTBOX_MARGIN = 0.5

# === TEXT GENERATION ===

def generate_prompt(template_str, values):
    return re.sub(r'\{\{(\w+)\}\}', lambda match: str(values.get(match.group(1), "")), template_str)

def request_claude_response(api_key, user_prompt):
    import urllib.request
    headers = {
        "x-api-key": api_key,
        "anthropic-version": "2023-06-01",
        "content-type": "application/json"
    }
    payload = {
        "model": "claude-3-5-sonnet-20241022",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": user_prompt}]
    }
    try:
        req = urllib.request.Request(
            url="https://api.anthropic.com/v1/messages",
            data=json.dumps(payload).encode("utf-8"),
            headers=headers,
            method="POST"
        )
        with urllib.request.urlopen(req) as response:
            result = json.loads(response.read())
            return result["content"][0]["text"]
    except Exception as error:
        print(f"[Claude Error] {error}")
        return None

# === DOCUMENT PARSING ===

def parse_docx_ordered(docx_path, image_dir="images"):
    from lxml import etree
    os.makedirs(image_dir, exist_ok=True)
    doc = Document(docx_path)
    ordered_items = []
    used_ids = set()
    img_count = 1

    for para in doc.paragraphs:
        if para.text.strip():
            ordered_items.append({"type": "text", "data": para.text.strip()})

        drawings = para._element.xpath(".//*[local-name()='drawing']")
        for draw in drawings:
            blip = draw.xpath(".//*[local-name()='blip']")
            if not blip:
                continue
            r_id = blip[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if r_id in used_ids:
                continue
            used_ids.add(r_id)
            image_bytes = doc.part.related_parts[r_id].blob
            filename = f"img_{img_count}_{uuid.uuid4().hex[:5]}.png"
            path = os.path.join(image_dir, filename)
            with open(path, "wb") as img_file:
                img_file.write(image_bytes)
            ordered_items.append({"type": "image", "path": path})
            img_count += 1
    return ordered_items

# === SLIDE CREATION HELPERS ===

def add_centered_image(slide, img_path, slide_w, slide_h):
    try:
        img = Image.open(img_path)
        w, h = img.size
        scale = min(MAX_IMAGE_WIDTH_INCHES / (w / 96), MAX_IMAGE_HEIGHT_INCHES / (h / 96), 1)
        w_scaled = Inches((w / 96) * scale)
        h_scaled = Inches((h / 96) * scale)
        left = (slide_w - w_scaled) / 2
        top = (slide_h - h_scaled) / 2
        slide.shapes.add_picture(img_path, left, top, width=w_scaled, height=h_scaled)
    except Exception as e:
        print(f"[Center Image Error] {e}")

def add_inline_image(slide, img_path, slide_w, slide_h, text_frame):
    try:
        img = Image.open(img_path)
        w, h = img.size
        scale = min(4.5 / (w / 96), 3.5 / (h / 96), 1)
        img_w = Inches((w / 96) * scale)
        img_h = Inches((h / 96) * scale)

        text_frame.width = slide_w - img_w - Inches(1.5)
        text_frame.left = Inches(0.5)
        text_frame.top = Inches(1.5)

        left = slide_w - img_w - Inches(0.5)
        top = (slide_h - img_h) / 2
        slide.shapes.add_picture(img_path, left, top, width=img_w, height=img_h)
    except Exception as e:
        print(f"[Inline Image Error] {e}")

def is_claude_fallback(text):
    phrases = [
        "not enough information", "i cannot", "provide more context", "please provide", "sorry", "missing content"
    ]
    return any(p in text.lower() for p in phrases)

# === GENERATE PRESENTATION ===

def build_presentation(api_key, docx_path, pptx_template, output_path):
    slides_content = parse_docx_ordered(docx_path)
    deck = Presentation(pptx_template)
    layout = deck.slide_layouts[1]
    slide_w, slide_h = deck.slide_width, deck.slide_height

    # clear existing slides
    for i in range(len(deck.slides) - 1, -1, -1):
        r_id = deck.slides._sldIdLst[i].rId
        deck.part.drop_rel(r_id)
        del deck.slides._sldIdLst[i]

    i = 0
    while i < len(slides_content):
        item = slides_content[i]

        # === TITLE SLIDE (first text + image) ===
        if i == 0 and item["type"] == "text":
            slide = deck.slides.add_slide(layout)
            slide.shapes.title.text = item["data"]

            if i + 1 < len(slides_content) and slides_content[i + 1]["type"] == "image":
                add_centered_image(slide, slides_content[i + 1]["path"], slide_w, slide_h)
                i += 1
            i += 1
            continue

        if item["type"] == "text":
            prompt = generate_prompt(GENERATE_SLIDE_CONTENT_TEMPLATE, {
                "topic": "",
                "contentSegment": item["data"]
            })
            response = request_claude_response(api_key, prompt)

            if not response or is_claude_fallback(response):
                i += 1
                continue

            lines = response.strip().splitlines()
            if lines[0].startswith("**"):
                title = lines[0].strip("* ")
                bullets = lines[1:]
            else:
                title = "Untitled"
                bullets = lines

            bullets = [b.strip("-• ") for b in bullets if b.strip()][:MAX_BULLETS_PER_SLIDE]
            if not bullets:
                i += 1
                continue

            slide = deck.slides.add_slide(layout)
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for b in bullets:
                tf.add_paragraph().text = b

            if i + 1 < len(slides_content) and slides_content[i + 1]["type"] == "image":
                add_inline_image(slide, slides_content[i + 1]["path"], slide_w, slide_h, slide.placeholders[1])
                i += 1

        i += 1

    deck.save(output_path)
    print(f"🎉 Presentation created at: {output_path}")

# === MAIN ENTRYPOINT ===

def main():
    load_dotenv()
    api_key = os.getenv("ANTHROPIC_API_KEY")

    base_path = os.path.dirname(os.path.abspath(__file__))
    input_doc = os.path.join(base_path, "document/doc.docx")
    pptx_template = os.path.join(base_path, "template/template.pptx")
    final_pptx = os.path.join(base_path, "output/output_presentation.pptx")

    build_presentation(api_key, input_doc, pptx_template, final_pptx)

if __name__ == "__main__":
    main()
